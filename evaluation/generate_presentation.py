import os
import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_slide_with_header(prs, title_text):
    """Creates a slide with a consistent header and styled title."""
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Add header title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Helvetica'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(44, 62, 80)
    
    # Accent line under title
    return slide

def add_bullet_points(slide, points, left, top, width, height, font_size=15):
    """Helper to populate bullet points into a textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = pt
        p.font.name = 'Helvetica'
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(68, 68, 68)
        p.space_after = Pt(10)
        p.level = 0
    return txBox

def main():
    base_dir = "/Users/archismanchoudhury/Desktop/research/Research Implementation"
    eval_dir = os.path.join(base_dir, "evaluation")
    docs_dir = os.path.join(base_dir, "docs")
    pptx_path = os.path.join(docs_dir, "Presentation.pptx")
    
    bench_df = pd.read_csv(os.path.join(eval_dir, "benchmark_results.csv"))
    perf_df = pd.read_csv(os.path.join(eval_dir, "performance_results.csv"))
    
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ------------------ SLIDE 1: Title ------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Fractional-Order Hyperchaotic Image Encryption\nUsing Sliding Mode Synchronization"
    p.font.name = 'Helvetica'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 119, 180)
    
    p2 = tf.add_paragraph()
    p2.text = "Design, Implementation, and Comparative Research Evaluation"
    p2.font.name = 'Helvetica'
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(44, 62, 80)
    p2.space_before = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "Presenter: Research Internship Presentation  |  Date: June 11, 2026"
    p3.font.name = 'Helvetica'
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(100, 100, 100)
    p3.space_before = Pt(40)
    
    # ------------------ SLIDE 2: Motivation ------------------
    slide2 = add_slide_with_header(prs, "Motivation")
    pts2 = [
        "Digital Image Security: Multimedia content transmission requires specialized high-speed ciphers due to high pixel redundancy and bulk data size.",
        "Limitations of Traditional Cryptography: AES/RSA block ciphers suffer from high computational latency and lack diffusion efficiency for pixel streams.",
        "Fractional Chaos Theory: Fractional-order systems (order q < 1) exhibit historical memory effects and infinite parameter spaces, expanding key complexity.",
        "Sliding Mode Control (SMC): Finite-time SMC ensures fast, noise-tolerant synchronization between decoupled Transmitter and Receiver attractors."
    ]
    add_bullet_points(slide2, pts2, Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    
    # ------------------ SLIDE 3: Problem Statement ------------------
    slide3 = add_slide_with_header(prs, "Problem Statement")
    pts3 = [
        "Active Synchronization Mismatch: High-precision chaotic ciphers suffer from tiny numerical floor errors (chattering) between drive and response states.",
        "Key Reconstruction Paradox: Direct quantization of chaotic variables causes mismatch bits at boundaries, resulting in decryption failure at the receiver.",
        "Common Workaround: Academic papers often bypass synchronization limitations by copying the exact transmitter key sequence directly to the receiver.",
        "Objective: Design and benchmark a parallel dual-mode secure pipeline:\n  - MODE = 'paper' (faithful direct trajectory quantization)\n  - MODE = 'robust' (entropy extraction + sign majority voting + SHA-256 + PRNG)"
    ]
    add_bullet_points(slide3, pts3, Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    
    # ------------------ SLIDE 4: Literature Review ------------------
    slide4 = add_slide_with_header(prs, "Literature Review")
    pts4 = [
        "Chaotic Cryptography (Sansom et al.): Pioneer ciphers utilized 1D logistic maps, vulnerable to phase space reconstruction and statistical attacks.",
        "Hyperchaotic Systems (Li et al.): Multi-scroll 6D systems generate multiple positive Lyapunov exponents, increasing randomness and keystream sensitivity.",
        "Fractional Dynamics (Petras et al.): Demonstrated memory effects where the derivative order 'q' serves as an additional key parameter.",
        "Sliding Mode Synchronization (Wang & Zhao): Explored projective SMC error correction, but ignored practical boundary quantization noise in secure communications."
    ]
    add_bullet_points(slide4, pts4, Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    
    # ------------------ SLIDE 5: Hyperchaotic System ------------------
    slide5 = add_slide_with_header(prs, "Hyperchaotic System Dynamics")
    pts5 = [
        "We implement a 6D fractional-order hyperchaotic system characterized by Caputo derivatives:",
        "  - x1' = a * (x2 - x1) + x4",
        "  - x2' = -x1 * x3 + b * x2",
        "  - x3' = x1 * x2 - c * x3",
        "  - x4' = d * x4 - x1 * x5",
        "  - x5' = e * x5 + x2 * x6",
        "  - x6' = -f * x6 + x3 * x4",
        "System parameters are initialized to: a = 0.5, b = 0.2, order q = 0.8.",
        "Adams-Bashforth-Moulton Solver: Evaluates memory truncation window (L = 2000) using the Short-Memory Principle to preserve fractional dynamics."
    ]
    add_bullet_points(slide5, pts5, Inches(1.0), Inches(1.5), Inches(11.33), Inches(5.5), font_size=14)
    
    # ------------------ SLIDE 6: Synchronization Framework ------------------
    slide6 = add_slide_with_header(prs, "Projective Sliding Mode Synchronization")
    pts6 = [
        "Drive-Response Setup: Transmitter runs Drive System (x), Receiver runs Response System (y).",
        "Error Vector: e(t) = y(t) - Gamma * x(t), with scaling matrix Gamma.",
        "Sliding Surface Design: Fractional integral sliding surface sigma(t) ensures finite-time convergence to zero error.",
        "Nonlinear Cancellation Controller: SMC control law eliminates system nonlinearities and injects regularized control input to minimize chattering.",
        "Stable Steady-State Convergence: Achieving an average error of ~1.54e-3 post-convergence at step 5001."
    ]
    add_bullet_points(slide6, pts6, Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    
    # ------------------ SLIDE 7: System Architecture ------------------
    slide7 = add_slide_with_header(prs, "System Architecture Flowchart")
    # Add diagram image
    img_path = os.path.join(docs_dir, "Architecture.png")
    if os.path.exists(img_path):
        slide7.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), Inches(10.33), Inches(5.2))
    else:
        add_bullet_points(slide7, ["[Architecture diagram image missing]"], Inches(1.0), Inches(2.0), Inches(11.0), Inches(4.0))
        
    # ------------------ SLIDE 8: Paper Mode Keystream ------------------
    slide8 = add_slide_with_header(prs, "Paper Mode: Direct Quantization")
    pts8 = [
        "Fidelity Focus: Keys are generated directly from trajectories without hashing, PRNGs, or majority voting.",
        "Noise-Free Subspace: Audit ranking identifies x5 as the most stable state (RMSE ~ 2.33e-4) and x2 as least stable (RMSE ~ 6.01e-3).",
        "Coarse Quantization Scheme: Uses only the most stable states (x1, x3, x5, x6) with noise-immune scaling factors:",
        "  - c_i = (q_1[i] * 27 + q_3[i] * 9 + q_5[i] * 3 + q_6[i]) mod 256",
        "  - q1, q3, q5 scale = 0.4 (values: 0, 1, 2); q6 scale = 0.1 (values: 0, 1).",
        "Guarantees 100% Key Parity: Eliminates boundary mismatches to enable perfect decryption, yielding 17 unique key values."
    ]
    add_bullet_points(slide8, pts8, Inches(1.0), Inches(1.5), Inches(11.33), Inches(5.5), font_size=14)
    
    # ------------------ SLIDE 9: Robust Mode Keystream ------------------
    slide9 = add_slide_with_header(prs, "Robust Mode: Cryptographic Seeding")
    pts9 = [
        "Engineering Focus: Employs sign majority voting to bridge synchronization fluctuations.",
        "Bit Extraction: Analyzes signs of states (x2 to x6) in blocks of size 32 (1 if positive, 0 if negative), extracting 10,240 bits.",
        "SHA-256 Seed Generation: Compresses raw bits into a uniform 256-bit seed, diffusing correlations and eliminating statistical bias.",
        "PRNG Keystream Expansion: Seeds a cryptographically secure NumPy PRNG to expand the keyspace uniformly over [0, 255].",
        "AVALANCHE Propagation: Guarantees high-entropy ciphers and total key sensitivity."
    ]
    add_bullet_points(slide9, pts9, Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    
    # ------------------ SLIDE 10: Security Results ------------------
    slide10 = add_slide_with_header(prs, "Security Metrics Results")
    # Add a table to slide 10
    rows = len(bench_df) + 1
    cols = 8
    table_shape = slide10.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8), Inches(12.33), Inches(4.5))
    table = table_shape.table
    
    # Column Headers
    headers = ["Image", "Mode", "Entropy", "NPCR %", "UACI %", "H Corr", "V Corr", "D Corr"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(44, 62, 80)
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Helvetica'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
    for r, row in bench_df.iterrows():
        vals = [
            str(row["Image"]), str(row["Mode"]),
            f"{row['Entropy']:.4f}", f"{row['NPCR']:.2f}%", f"{row['UACI']:.2f}%",
            f"{row['H_Corr']:.4f}", f"{row['V_Corr']:.4f}", f"{row['D_Corr']:.4f}"
        ]
        for c, val in enumerate(vals):
            cell = table.cell(r + 1, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Helvetica'
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER
            
    # ------------------ SLIDE 11: Performance Results ------------------
    slide11 = add_slide_with_header(prs, "Runtime Performance & Throughput")
    rows = len(perf_df) + 1
    cols = 6
    table_shape2 = slide11.shapes.add_table(rows, cols, Inches(1.0), Inches(2.0), Inches(11.33), Inches(2.5))
    table2 = table_shape2.table
    
    headers2 = ["Mode", "Enc Time (s)", "Dec Time (s)", "Enc Throughput (MB/s)", "Dec Throughput (MB/s)", "Peak Memory (MB)"]
    for c, h in enumerate(headers2):
        cell = table2.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(44, 62, 80)
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Helvetica'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
    for r, row in perf_df.iterrows():
        vals = [
            str(row["Mode"]).upper(),
            f"{row['Encryption_Time_Avg_s']:.5f}",
            f"{row['Decryption_Time_Avg_s']:.5f}",
            f"{row['Encryption_Throughput_MBs']:.2f}",
            f"{row['Decryption_Throughput_MBs']:.2f}",
            f"{row['Peak_Process_Memory_MB']:.2f}"
        ]
        for c, val in enumerate(vals):
            cell = table2.cell(r + 1, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Helvetica'
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER
            
    # Add brief bullet points under table
    pts11 = [
        "Synchronization overhead is a constant ~13.77s for 70,536 integration steps (equal for both modes).",
        "Encryption throughput is equivalent (~1.92-1.98 MB/s), dominated by modulo-256 diffusion loop execution."
    ]
    add_bullet_points(slide11, pts11, Inches(1.0), Inches(5.0), Inches(11.33), Inches(2.0))
    
    # ------------------ SLIDE 12: Robustness Under Noise ------------------
    slide12 = add_slide_with_header(prs, "Robustness Under Synchronization Noise")
    pts12 = [
        "Robustness Benchmark: Add Gaussian noise sigma to response trajectory y(t).",
        "Robust Mode Cliff-Edge Effect: Sign majority voting filter tolerates noise up to sigma = 5e-3. Beyond this threshold, a single flipped bit changes the SHA-256 seed, dropping key parity instantly to 0.39% and failing decryption.",
        "Paper Mode Gradual Degradation: Trajectory noise causes localized grid boundary crossings. Key parity degrades slowly (99.91% at 1e-6, 99.49% at 1e-2).",
        "Symmetric Diffusion Limitation: Because of carry propagation in modulo-256 diffusion, even a 0.09% key difference completely corrupts image decryption. Thus, Paper Mode decryption still fails under active noise."
    ]
    add_bullet_points(slide12, pts12, Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    
    # ------------------ SLIDE 13: Comparative Analysis ------------------
    slide13 = add_slide_with_header(prs, "Comparative Analysis & Trade-offs")
    pts13 = [
        "Cryptographic Quality: Robust Mode guarantees uniform keystream distribution and high entropy (7.98), whereas Paper Mode is restricted to 17 unique key values (7.96 entropy).",
        "Keyspace Security: Robust Mode has $2^256$ seed keyspace, ensuring brute-force resistance. Paper Mode keyspace is theoretically large ($2^376$) but practically weak due to keystream structure.",
        "Fidelity: Paper Mode stays 100% faithful to trajectory direct-quantization, bypassing all non-physical hashes, PRNGs, and seed generators.",
        "Deployment Recommendation: Robust Mode is recommended for real-world deployments due to high cryptographic margin and low-pass filter noise resilience."
    ]
    add_bullet_points(slide13, pts13, Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    
    # ------------------ SLIDE 14: Future Work ------------------
    slide14 = add_slide_with_header(prs, "Future Research Directions")
    pts14 = [
        "Color Image Encryption: Extend the modular permutation and diffusion framework to RGB channels.",
        "GPU Solver Acceleration: Port the ABM-PC solver and SMC calculations to CUDA/OpenCL to enable real-time execution.",
        "Hardware Co-Processor: Design and implement the fractional synchronization controller on FPGA.",
        "Adaptive SMC Controller: Research adaptive sliding mode control to dynamically adjust gains under active channel noise."
    ]
    add_bullet_points(slide14, pts14, Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    
    # ------------------ SLIDE 15: Q&A / Conclusion ------------------
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    tb_qa = slide15.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf_qa = tb_qa.text_frame
    tf_qa.word_wrap = True
    p = tf_qa.paragraphs[0]
    p.text = "Thank You"
    p.font.name = 'Helvetica'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 119, 180)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf_qa.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.font.name = 'Helvetica'
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(44, 62, 80)
    p2.space_before = Pt(20)
    p2.alignment = PP_ALIGN.CENTER
    
    prs.save(pptx_path)
    print(f"[+] Successfully generated PowerPoint slide deck: {pptx_path}")

if __name__ == "__main__":
    main()
